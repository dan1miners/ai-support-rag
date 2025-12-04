# Add these two lines at the very top
from dotenv import load_dotenv
load_dotenv() # Loads variables from .env file

import os
import io
import uuid
from typing import List, Optional

import supabase
from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
import openai
from postgrest.exceptions import APIError
from fastapi.middleware.cors import CORSMiddleware

# --- CONFIGURATION ---
# Now os.getenv() will correctly find the variables loaded from .env
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")




# Check if the credentials are loaded
if not all([SUPABASE_URL, SUPABASE_KEY, OPENAI_API_KEY]):
    raise ValueError("Missing one or more environment variables: SUPABASE_URL, SUPABASE_KEY, OPENAI_API_KEY")

# Supabase client
supabase_client = supabase.create_client(SUPABASE_URL, SUPABASE_KEY)

# OpenAI client
openai_client = openai.OpenAI(api_key=OPENAI_API_KEY)

# FastAPI app
app = FastAPI(title="RAG Backend API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allows all origins for development. For production, specify your Streamlit app's URL.
    allow_credentials=True,
    allow_methods=["*"],  # Allows all methods
    allow_headers=["*"],  # Allows all headers
)


# --- DATA MODELS ---
class ChatRequest(BaseModel):
    query: str
    version_id: Optional[int] = None

class ChatResponse(BaseModel):
    answer: str
    sources: List[dict]

# --- HELPER FUNCTIONS ---
def process_file(file: UploadFile) -> str:
    """Extracts text content from an uploaded file."""
    content = file.file.read()
    
    if file.content_type == "text/plain":
        return content.decode("utf-8")
    elif file.content_type == "application/pdf":
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(content))
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    elif file.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        import docx
        doc = docx.Document(io.BytesIO(content))
        return "\n".join([para.text for para in doc.paragraphs])
    elif file.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        import pptx
        prs = pptx.Presentation(io.BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif file.content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        import openpyxl
        workbook = openpyxl.load_workbook(io.BytesIO(content))
        text = ""
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text:
                    text += row_text + "\n"
        return text
    else:
        raise ValueError(f"Unsupported file type: {file.content_type}")

def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    """Splits text into smaller chunks."""
    # The import statement is now from the new, separate package
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=chunk_overlap
    )
    return text_splitter.split_text(text)

def create_embedding(text: str) -> List[float]:
    """Creates an embedding for a given text using OpenAI."""
    response = openai_client.embeddings.create(
        model="text-embedding-3-small", input=text
    )
    return response.data[0].embedding

# --- API ENDPOINTS ---

# --- API ENDPOINTS ---

@app.post("/upload/")
async def upload_and_process_files(
    files: List[UploadFile] = File(...),
    create_new_version: bool = Form(False),  # Default: add to active version
    target_version_id: Optional[int] = Form(None)  # Optional: specific version
):
    """
    Uploads files with flexible version handling.
    
    Options:
    1. create_new_version=True: Create new version (default: active if none exists)
    2. target_version_id provided: Add to specific version
    3. Neither: Add to active version (create if none)
    """
    if not files:
        raise HTTPException(status_code=400, detail="No files provided.")

    # Logic to determine target version
    if target_version_id:
        # Add to specific version
        version_check = supabase_client.table("versions").select("*").eq("id", target_version_id).execute()
        if not version_check.data:
            raise HTTPException(status_code=404, detail=f"Version {target_version_id} not found")
        
        target_version = version_check.data[0]
        target_version_id = target_version['id']
        version_name = target_version['name']
        is_new_version = False
        # Keep the version's active status as is
        is_active = target_version.get('is_active', False)
        
    elif create_new_version:
        # Create new version
        try:
            active_versions = supabase_client.table("versions").select("id").eq("is_active", True).execute()
            has_active_version = len(active_versions.data) > 0
        except Exception as e:
            print(f"Error checking active versions: {e}")
            has_active_version = False

        version_name = f"v_{uuid.uuid4().hex[:6]}"
        version_response = supabase_client.table("versions").insert({
            "name": version_name,
            "is_active": not has_active_version  # Active if no active version exists
        }).execute()
        
        if not version_response.data:
            raise HTTPException(status_code=500, detail="Failed to create new version.")
        
        target_version_id = version_response.data[0]['id']
        is_new_version = True
        is_active = not has_active_version
        
    else:
        # Add to active version (or create if none)
        try:
            version_response = supabase_client.table("versions").select("*").eq("is_active", True).execute()
            
            if version_response.data:
                # Use existing active version
                active_version = version_response.data[0]
                target_version_id = active_version['id']
                version_name = active_version['name']
                is_new_version = False
                is_active = True
            else:
                # Create a new active version
                version_name = f"v_{uuid.uuid4().hex[:6]}"
                version_response = supabase_client.table("versions").insert({
                    "name": version_name,
                    "is_active": True
                }).execute()
                
                if not version_response.data:
                    raise HTTPException(status_code=500, detail="Failed to create new version.")
                
                target_version_id = version_response.data[0]['id']
                is_new_version = True
                is_active = True
                
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error checking versions: {str(e)}")

    processed_chunks = 0

    # Process files
    for file in files:
        try:
            text_content = process_file(file)
            if not text_content:
                continue

            chunks = chunk_text(text_content)

            for chunk in chunks:
                embedding = create_embedding(chunk)
                supabase_client.table("documents").insert({
                    "content": chunk,
                    "metadata": {"source": file.filename},
                    "embedding": embedding,
                    "version_id": target_version_id
                }).execute()
                processed_chunks += 1
        
        except Exception as e:
            print(f"Error processing file {file.filename}: {e}")
            continue

    return {
        "message": "Files processed successfully.",
        "version_name": version_name,
        "version_id": target_version_id,
        "chunks_created": processed_chunks,
        "is_new_version": is_new_version,
        "is_active": is_active
    }

@app.get("/versions/")
def get_all_versions():
    """Retrieves a list of all knowledge base versions with document counts."""
    try:
        # Get all versions
        response = supabase_client.table("versions").select("*").order("created_at", desc=True).execute()
        
        # For each version, get the document count
        versions_with_counts = []
        for version in response.data:
            # Count documents for this version
            doc_count_response = supabase_client.table("documents")\
                .select("id", count="exact")\
                .eq("version_id", version['id'])\
                .execute()
            
            versions_with_counts.append({
                "id": version['id'],
                "version_name": version.get('name', 'Unnamed Version'),  # Map 'name' to 'version_name'
                "name": version.get('name', 'Unnamed Version'),  # Keep original for compatibility
                "is_active": version.get('is_active', False),
                "created_at": version.get('created_at'),
                "document_count": doc_count_response.count or 0
            })
        
        return versions_with_counts
        
    except Exception as e:
        print(f"Error fetching versions: {e}")
        raise HTTPException(status_code=500, detail=f"Error fetching versions: {str(e)}")

@app.post("/chat/")
async def chat_with_rag(request: ChatRequest):
    """
    Handles a user query, finds relevant documents, and generates an answer.
    """
    print(f"Received chat request: {request.query}")  # Add this for debugging
    
    # 1. Find the currently active version
    try:
        version_response = supabase_client.table("versions").select("id").eq("is_active", True).single().execute()
        active_version_id = version_response.data['id']
        print(f"Active version ID: {active_version_id}")  # Debug
    except Exception as e:
        print(f"Error finding active version: {e}")  # Debug
        # Return a simple streaming response for the error
        def error_stream():
            yield "No active knowledge base version found. Please upload a document first."
        
        return StreamingResponse(error_stream(), media_type="text/plain")

    # 2. Create embedding for the user's query
    query_embedding = create_embedding(request.query)

    # 3. Search for similar document chunks
    try:
        match_response = supabase_client.rpc(
            "match_documents",
            {
                "query_embedding": query_embedding,
                "match_count": 5,
                "version_filter": active_version_id
            }
        ).execute()
        print(f"Found {len(match_response.data) if match_response.data else 0} matching documents")  # Debug
    except Exception as e:
        print(f"Search error: {e}")  # Debug
        def search_error_stream():
            yield f"Search error: {str(e)}"
        
        return StreamingResponse(search_error_stream(), media_type="text/plain")

    if not match_response.data:
        def no_results_stream():
            yield "I couldn't find any relevant information in the knowledge base."
        
        return StreamingResponse(no_results_stream(), media_type="text/plain")

    # 4. Format context for the LLM
    context = "\n\n".join([doc["content"] for doc in match_response.data])

    # 5. Generate answer using OpenAI
    system_prompt = """You are OneMiners AI Assistant, the official helper for the OneMiners community. 
        Your mission is to provide accurate information with warmth and empathy.

        **Response Rules:**
        1. Base answers ONLY on the provided context
        2. If information is missing, respond kindly: "I don't have that information in my current knowledge base"
        3. Always be supportive and encouraging
        4. Use a friendly, professional tone
        5. End responses with an offer to help with other questions

        **Example Response Patterns:**
        - When you know: "Based on the available information, I can help with that! [Answer]"
        - When you don't know: "I want to help you with that, but I don't have that specific information in my current knowledge base. Would you like me to help with something else from the documents I do have access to?"
        """

    user_prompt = f"""**Knowledge Base Context:**
        {context}

        **User's Question:**
        "{request.query}"

        **Your Response (as OneMiners AI Assistant):**
        """
    try:
        stream = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            stream=True,
            temperature=0.7,  
        )
    except Exception as e:
        print(f"OpenAI error: {e}")  # Debug
        def openai_error_stream():
            yield f"OpenAI error: {str(e)}"
        
        return StreamingResponse(openai_error_stream(), media_type="text/plain")

    def stream_generator():
        full_response = ""
        for chunk in stream:
            content = chunk.choices[0].delta.content or ""
            full_response += content
            yield content
    
    return StreamingResponse(stream_generator(), media_type="text/event-stream")


@app.post("/chat/non-streaming/", response_model=ChatResponse)
async def chat_with_rag_non_streaming(request: ChatRequest):
    """
    Non-streaming version of chat endpoint for debugging and clients that don't support streaming.
    """
    # 1. Find the currently active version
    try:
        version_response = supabase_client.table("versions").select("id").eq("is_active", True).single().execute()
        active_version_id = version_response.data['id']
    except APIError as e:
        if e.code == 'PGRST116':
            raise HTTPException(
                status_code=404, 
                detail="No active knowledge base version found. Please upload a document and set a version as active."
            )
        else:
            raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Unexpected error: {str(e)}")

    # 2. Create embedding for the user's query
    query_embedding = create_embedding(request.query)

    # 3. Search for similar document chunks
    try:
        match_response = supabase_client.rpc(
            "match_documents",
            {
                "query_embedding": query_embedding,
                "match_count": 5,
                "version_filter": active_version_id
            }
        ).execute()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Search error: {str(e)}")

    if not match_response.data:
        return ChatResponse(
            answer="I couldn't find any relevant information in the knowledge base.", 
            sources=[]
        )

    # 4. Format context for the LLM
    context = "\n\n".join([doc["content"] for doc in match_response.data])
    sources = [{"content": doc["content"][:200], "metadata": doc.get("metadata", {})} for doc in match_response.data]

    # 5. Generate answer using OpenAI
    system_prompt = "You are a helpful assistant. Answer the user's question based only on the provided context. If the context doesn't contain the answer, say you don't know."
    user_prompt = f"Context:\n{context}\n\nQuestion: {request.query}"

    try:
        completion = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            stream=False,
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"OpenAI error: {str(e)}")
    
    answer = completion.choices[0].message.content
    
    return ChatResponse(answer=answer, sources=sources)


# Add this endpoint right after your /versions/ endpoint

@app.post("/versions/{version_id}/switch")
def switch_active_version(version_id: int):
    """Sets a specific version as the active one for queries."""
    try:
        # First, check if the version exists
        version_check = supabase_client.table("versions").select("*").eq("id", version_id).execute()
        if not version_check.data:
            raise HTTPException(status_code=404, detail="Version not found.")
        
        # 1. Deactivate all versions
        deactivate_response = supabase_client.table("versions").update({"is_active": False}).eq("is_active", True).execute()
        print(f"Deactivated {len(deactivate_response.data) if deactivate_response.data else 0} active versions")
        
        # 2. Activate the selected version
        activate_response = supabase_client.table("versions").update({"is_active": True}).eq("id", version_id).execute()
        
        if not activate_response.data:
            raise HTTPException(status_code=500, detail="Failed to activate version.")
        
        version_name = activate_response.data[0].get('name', f"Version {version_id}")
        
        return {
            "message": f"Successfully switched to version: {version_name}",
            "version_id": version_id,
            "version_name": version_name
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error switching version: {e}")
        raise HTTPException(status_code=500, detail=f"Error switching version: {str(e)}")
    
@app.get("/", tags=["Root"])
async def read_root():
    return {
        "message": "Welcome to the RAG Backend API. The service is running.",
        "docs": "/docs",
        "redoc": "/redoc"
    }
    
    
    # Add this at the very end of your main.py file
@app.get("/debug/routes")
def debug_routes():
    """
    A temporary endpoint to print all registered routes.
    """
    routes = []
    for route in app.routes:
        routes.append({
            "path": route.path,
            "methods": [method.upper() for method in route.methods]
        })
    return routes

# Add this endpoint in your main.py (FastAPI backend)
@app.delete("/versions/{version_id}")
def delete_version(version_id: int):
    """
    Deletes a version and all its associated documents.
    """
    try:
        # 1. First check if version exists
        version_check = supabase_client.table("versions").select("*").eq("id", version_id).execute()
        if not version_check.data:
            raise HTTPException(status_code=404, detail=f"Version {version_id} not found.")
        
        version_data = version_check.data[0]
        version_name = version_data.get('name', f"Version {version_id}")
        is_active = version_data.get('is_active', False)
        
        # 2. Check if trying to delete active version (optional restriction)
        if is_active:
            return {
                "warning": "Cannot delete the active version. Please switch to another version first.",
                "version_id": version_id,
                "version_name": version_name,
                "status": "active"
            }
        
        # 3. First, delete all documents associated with this version
        try:
            # Get count of documents to delete
            doc_count_response = supabase_client.table("documents")\
                .select("id", count="exact")\
                .eq("version_id", version_id)\
                .execute()
            
            doc_count = doc_count_response.count or 0
            
            # Delete documents
            delete_docs_response = supabase_client.table("documents")\
                .delete()\
                .eq("version_id", version_id)\
                .execute()
            
            print(f"Deleted {doc_count} documents for version {version_id}")
            
        except Exception as e:
            print(f"Error deleting documents: {e}")
            # Continue to try deleting the version anyway
        
        # 4. Delete the version record
        delete_version_response = supabase_client.table("versions")\
            .delete()\
            .eq("id", version_id)\
            .execute()
        
        if not delete_version_response.data:
            raise HTTPException(status_code=500, detail=f"Failed to delete version {version_id}")
        
        return {
            "message": f"Version '{version_name}' deleted successfully.",
            "version_id": version_id,
            "version_name": version_name,
            "documents_deleted": doc_count,
            "status": "deleted"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error deleting version: {e}")
        raise HTTPException(status_code=500, detail=f"Error deleting version: {str(e)}")